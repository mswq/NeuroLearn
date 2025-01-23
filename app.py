from flask import Flask, request, render_template, redirect, url_for, flash
from flask_sqlalchemy import SQLAlchemy
from flask_bcrypt import Bcrypt
from flask_login import UserMixin, LoginManager, login_user, logout_user, login_required, current_user
from werkzeug.utils import secure_filename
import os
from pptx import Presentation
import openai
from dotenv import load_dotenv

load_dotenv()

app = Flask(__name__)
app.config["UPLOAD_FOLDER"] = "./uploads"
app.config["SQLALCHEMY_DATABASE_URI"] = "sqlite:///users.db"
app.config["SQLALCHEMY_TRACK_MODIFICATIONS"] = False
app.secret_key = "secret_key"

db = SQLAlchemy(app)
bcrypt = Bcrypt(app)
login_manager = LoginManager(app)
login_manager.login_view = "login"

openai.api_key = os.getenv("OPENAI_API_KEY")

class User(db.Model, UserMixin):
    id = db.Column(db.Integer, primary_key=True)
    username = db.Column(db.String(150), unique=True, nullable=False)
    password = db.Column(db.String(150), nullable=False)

@app.route("/")
def index():
    return render_template("index.html")

@app.route("/get_started")
def get_started():
    return render_template("get_started.html")

@app.route("/about")
def about():
    return render_template("about.html")

@app.route('/study_plan')
def study_plan():
    return render_template("study_plan.html")

@app.route('/upload', methods=['POST'])
def upload_file():
    try:
        neurotype = request.form.get('neurotype')
        learning_style = request.form.get('learning_style')
        
        if 'file' not in request.files:
            flash("No file uploaded.", "error")
            return redirect(url_for('get_started'))
        
        file = request.files['file']

        if file.filename == '':
            flash("No file selected.", "error")
            return redirect(url_for('get_started'))
        if not file.filename.endswith('.pptx'):
            flash("Invalid file type. Only .pptx files are allowed.", "error")
            return redirect(url_for('get_started'))

        safe_filename = secure_filename(file.filename)
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], safe_filename)
        os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
        file.save(filepath)

        study_plan = generate_study_plan_with_openai(neurotype, learning_style, filepath)

        return render_template(
            'study_plan.html',
            neurotype=neurotype,
            learning_style=learning_style,
            study_plan=study_plan
        )

    except Exception as e:
        print(f"Error: {e}")
        return "An error occurred. Please try again.", 500
    
def generate_study_plan_with_openai(neurotype, learning_style, filepath):
    try:
        content = extract_pptx_content(filepath)
        print(f"Extracted Content: {content}") 

        prompt = construct_prompt(neurotype, learning_style, content)
        print(f"Generated Prompt: {prompt}") 

        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You are an expert educator specializing in personalized learning plans."},
                {"role": "user", "content": prompt},
            ],
            max_tokens=500,
            temperature=0.7,
        )
        return response["choices"][0]["message"]["content"].strip()

    except Exception as e:
        import traceback
        print(f"OpenAI API Error:\n{traceback.format_exc()}")
        return "An error occurred while generating the study plan."

def extract_pptx_content(filepath):
    presentation = Presentation(filepath)
    slides_content = []

    for slide in presentation.slides:
        slide_data = {"title": None, "content": []}
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.strip()
                if not slide_data["title"]:
                    slide_data["title"] = text  
                else:
                    slide_data["content"].append(text)
        slides_content.append(slide_data)

    formatted_content = "\n".join(
        f"{slide['title']}:\n{' '.join(slide['content'])}" for slide in slides_content
    )
    return formatted_content

def construct_prompt(neurotype, learning_style, content):
    prompt = f"""
    Create a personalized study plan for a student based on the following details:
    - Neurodiverse Type: {neurotype}
    - Preferred Learning Style: {learning_style}
    - Content: {content}

    Provide actionable steps, tailored to the student's needs, in a clear and concise format.
    """
    return prompt

@app.route('/signup', methods=['GET', 'POST'])
def signup():
    if request.method =='POST':
        username = request.form['username']
        password = request.form['password']

        existing_user = User.query.filter_by(username=username).first()
        if existing_user:
            flash('Username already exists.', 'error')
            return redirect(url_for('signup'))
        
        hashed_password = bcrypt.generate_password_hash(password).decode('utf-8')
        new_user = User(username=username, password=hashed_password)
        db.session.add(new_user)
        db.session.commit()

        flash('Account created successfully. Please log in.', 'success')
        return redirect(url_for('login'))
    return render_template('signup.html')

@login_manager.user_loader  
def load_user(user_id):
    return User.query.get(int(user_id))

@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        username = request.form['username']
        password = request.form['password']

        user = User.query.filter_by(username=username).first()
        if user and bcrypt.check_password_hash(user.password, password):
            login_user(user)
            flash("Logged in successfully", 'success')
            return redirect(url_for('dashboard'))
        else:
            flash('Invalid username or password. Please try again.', 'error')
    return render_template('login.html')

@app.route('/dashboard')
@login_required
def dashboard():
    return f"Welcome to your dashboard, {current_user.username}!"

if __name__ == "__main__":
    with app.app_context():
        db.create_all()
    app.run(debug=True)
